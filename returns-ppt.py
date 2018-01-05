import time
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt


PIC_PATH = "/home/sma-analytics/Data/Returns/"
LEFT = Inches(0)
TOP = Inches(1)
###Update Title Slide
prez = Presentation(PIC_PATH + 'Returns.pptx')
i = 0
date = time.strftime("%b %d, %Y")
pres_name = PIC_PATH + "Returns.pptx"
prez.slides[0].shapes[4].text = date
p = prez.slides[0].shapes[4].text_frame.paragraphs[0]
p.font.size = Pt(18)
p.font.name = 'Segoe'
p.font.bold = True


###YTD-Long Short (Slide 3)
#To Do:
    #Graph
    #Values Chart

cur_slide = prez.slides[2]

#cur_slide.shapes.pop(1)
#cur_slide.shapes.pop(1)

cur_slide.shapes.add_picture(PIC_PATH + 'YTD_LS.png', LEFT, TOP)


###Rolling 1 Year - Long Short (Slide 4)
#To Do:
    #Graph
    #Values Chart

cur_slide = prez.slides[3]

#cur_slide.shapes.pop(1)
#cur_slide.shapes.pop(1)

cur_slide.shapes.add_picture(PIC_PATH + 'Rolling_LS.png', LEFT, TOP)


###Rolling 1 Year - Long Only (Slide 5)
#To Do:
    #Graph
    #Values Chart

cur_slide = prez.slides[4]

#cur_slide.shapes.pop(1)
#cur_slide.shapes.pop(1)

cur_slide.shapes.add_picture(PIC_PATH + 'Rolling_Long.png', LEFT, TOP)


###2016 to Current - Long Short (Slide 6)
#To Do:
    #Graph
    #Values Chart

cur_slide = prez.slides[5]

#cur_slide.shapes.pop(1)
#cur_slide.shapes.pop(1)

cur_slide.shapes.add_picture(PIC_PATH + '2016plus_LS.png', LEFT, TOP)


###2016 to Current - Long Only (Slide 7)
#To Do:
    #Graph
    #Values Chart

cur_slide = prez.slides[6]

#cur_slide.shapes.pop(1)
#cur_slide.shapes.pop(1)

cur_slide.shapes.add_picture(PIC_PATH + '2016plus_Long.png', LEFT, TOP)


###Full History - Long Short (Slide 8)
#To Do:
    #Graph
    #Values Chart

cur_slide = prez.slides[7]

#cur_slide.shapes.pop(1)
#cur_slide.shapes.pop(1)

cur_slide.shapes.add_picture(PIC_PATH + 'FullHistory_LS.png', LEFT, TOP)



###Full History - Long Only (Slide 9)
#To Do:
    #Graph
    #Values Chart

cur_slide = prez.slides[8]

#cur_slide.shapes.pop(1)
#cur_slide.shapes.pop(1)

cur_slide.shapes.add_picture(PIC_PATH + 'FullHistory_Long.png', LEFT, TOP)




prez.save(pres_name)
    
    
